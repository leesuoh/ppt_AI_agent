import os, subprocess
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from utils import safe_run, save_state_log
from PIL import Image as PILImage
import io, base64

MEDIA_DIR = "output"

def clean_text(text):
    return text.replace("\n", " ").strip()

def split_sents(text):
    return [s.strip() for s in text.split(".") if s.strip()]


@safe_run
def node_parse_ppt(state: dict) -> dict:
    pptx_path = state["pptx_path"]
    slide_index = state.get("slide_index", 0)

    # LibreOffice로 PDF 변환
    subprocess.run([
        "libreoffice", "--headless", "--convert-to", "pdf", pptx_path,
        "--outdir", MEDIA_DIR
    ], check=True)
    ppt_pdf_name = os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf"
    pdf_path = os.path.join(MEDIA_DIR, ppt_pdf_name)

    # PDF → PNG 변환 (Poppler 사용)
    slide_png_pattern = os.path.join(MEDIA_DIR, f"slide{slide_index}_full")
    subprocess.run([
        "pdftoppm", "-f", str(slide_index+1), "-l", str(slide_index+1),
        "-png", pdf_path, slide_png_pattern
    ], check=True)
    slide_image_path = slide_png_pattern + "-1.png"

    # PPTX 열어서 텍스트/테이블/이미지 추출
    ppt = Presentation(pptx_path)
    slide = ppt.slides[slide_index]

    texts, tables, images = [], [], []

    for i, sh in enumerate(slide.shapes):
        if sh.has_text_frame:
            txt = "\n".join(p.text for p in sh.text_frame.paragraphs)
            texts.append(txt)
        if sh.shape_type == MSO_SHAPE_TYPE.TABLE:
            tbl = [[clean_text(c.text) for c in r.cells] for r in sh.table.rows]
            tables.append(tbl)
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            ext = sh.image.ext
            path = os.path.join(MEDIA_DIR, f"slide{slide_index}_img{i}.{ext}")
            with open(path, "wb") as f:
                f.write(sh.image.blob)
            images.append(path)

    # 텍스트 전처리
    txt_c = clean_text(" ".join(texts))
    txt_ss = split_sents(txt_c)

    state["texts"] = txt_ss
    state["tables"] = tables
    state["images"] = images
    state["slide_image"] = [slide_image_path]

    save_state_log(state, "parse_ppt")
    return state
